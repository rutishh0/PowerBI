"""
Rolls-Royce Civil Aerospace — Data Visualizer Server
=====================================================
Flask API backend serving the premium dashboard frontend.
Handles file upload, universal Excel parsing, and PDF export.
Supports: SOA, Invoice List, Opportunity Tracker, Shop Visit History, SVRG Master.

Usage:
    python server.py
    Then open http://localhost:5000 in your browser.
"""

import io
import os
import json
import uuid
import math
import base64
import time
import threading
import concurrent.futures

import pandas as pd
from flask import Flask, request, jsonify, send_file, send_from_directory, session, redirect
from functools import wraps
from flask_cors import CORS

# Universal parser — handles SOA, INVOICE_LIST, OPPORTUNITY_TRACKER, SHOP_VISIT, SVRG_MASTER
from parser import parse_file
# Keep old parser for PDF export backward compat
from parser import parse_soa_workbook, serialize_parsed_data, aging_bucket, fmt_currency, AGING_ORDER, AGING_COLORS
from pdf_export import generate_pdf_report
from ai_chat import build_system_prompt, call_openrouter
from storage import init_db, save_file_to_db, kv_get, kv_set, r2_get_text, r2_put_text
from storage import (
    download_from_r2, delete_from_r2,
    generate_r2_key, create_multipart_upload, upload_part,
    complete_multipart_upload, abort_multipart_upload, R2_PUBLIC_URL,
)
import datetime as dt

# Tracks active multipart uploads: { session_upload_id: { r2_key, r2_upload_id, parts[], filename, total, file_size } }
_multipart_sessions = {}

def _sanitize_for_json(obj):
    """Recursively sanitize parsed data for JSON serialization.
    Converts NaN/Infinity to None, datetime to str, pandas types to native Python."""
    if obj is None:
        return None
    if isinstance(obj, float):
        if math.isnan(obj) or math.isinf(obj):
            return None
        return obj
    if isinstance(obj, (int, bool, str)):
        return obj
    if isinstance(obj, (dt.datetime, dt.date)):
        return obj.isoformat()
    if isinstance(obj, dict):
        return {k: _sanitize_for_json(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple)):
        return [_sanitize_for_json(v) for v in obj]
    # Handle pandas types, numpy types
    try:
        import numpy as np
        if isinstance(obj, (np.integer,)):
            return int(obj)
        if isinstance(obj, (np.floating,)):
            v = float(obj)
            return None if math.isnan(v) or math.isinf(v) else v
        if isinstance(obj, np.ndarray):
            return [_sanitize_for_json(v) for v in obj.tolist()]
        if isinstance(obj, np.bool_):
            return bool(obj)
    except ImportError:
        pass
    try:
        import pandas as pd
        if isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        if pd.isna(obj):
            return None
    except (ImportError, TypeError, ValueError):
        pass
    # Fallback: convert to string
    return str(obj)

# ─────────────────────────────────────────────────────────────
# APP SETUP
# ─────────────────────────────────────────────────────────────

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "rr-soa-dashboard-" + uuid.uuid4().hex[:8])
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50MB per request (R2 chunks are ~11MB each)

# Session cookie config — tuned so the Next.js frontend can reach this API
# either same-origin (via Next rewrites in dev/prod) or cross-origin with
# credentials. Browsers require Secure=True for SameSite=None; localhost is
# treated as secure by Chrome/Firefox so this works in dev too.
_SAMESITE = os.environ.get("SESSION_COOKIE_SAMESITE", "Lax")  # "None" for cross-origin direct
app.config["SESSION_COOKIE_SAMESITE"] = _SAMESITE
app.config["SESSION_COOKIE_SECURE"] = _SAMESITE.lower() == "none" or os.environ.get(
    "SESSION_COOKIE_SECURE", ""
).lower() in ("1", "true", "yes")
app.config["SESSION_COOKIE_HTTPONLY"] = True

# CORS — allow the Next frontend origin(s). In production set CORS_ORIGINS
# to a comma-separated list (e.g. "https://dashboard.example.com").
_origins_env = os.environ.get(
    "CORS_ORIGINS",
    "http://localhost:3000,http://127.0.0.1:3000",
)
_cors_origins = [o.strip() for o in _origins_env.split(",") if o.strip()]
CORS(app, origins=_cors_origins, supports_credentials=True)

# Initialize DB tables on import (needed for gunicorn which skips __main__)
init_db()

# ── Feature Flags ──────────────────────────────────────────
# Set ENABLE_EXTRA_FEATURES to True to re-enable AI chat, file vault,
# comparison mode, and the secret admin chat button.
ENABLE_EXTRA_FEATURES = True

FEATURE_FLAGS = {
    "show_ai":          ENABLE_EXTRA_FEATURES,
    "show_files":       ENABLE_EXTRA_FEATURES,
    "show_compare":     ENABLE_EXTRA_FEATURES,
    "show_secret_chat": ENABLE_EXTRA_FEATURES,
}

# In-memory store for parsed data (keyed by session ID)
# In production, use Redis or similar
_parsed_store = {}

# In-memory store for chat history (keyed by session ID)
_chat_history = {}

# In-memory store for AI report generation jobs (keyed by job_id). Per-process,
# like _parsed_store — the service must run effectively single-worker (gunicorn
# --workers 1 --threads N). Jobs hold the finished PDF until downloaded.
_ai_report_jobs = {}
_ai_report_lock = threading.Lock()
_AI_JOB_TTL = 1800  # seconds; prune finished/stale jobs older than this


def _prune_ai_jobs():
    now = time.time()
    with _ai_report_lock:
        for jid in [k for k, v in _ai_report_jobs.items()
                    if now - v.get("created", now) > _AI_JOB_TTL]:
            _ai_report_jobs.pop(jid, None)


def _get_session_id():
    """Get or create a session ID."""
    if "sid" not in session:
        session["sid"] = uuid.uuid4().hex
    return session["sid"]


def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get("authenticated"):
            # All login_required routes are /api/*; JSON 401 is what the
            # frontend expects (ApiError handler in lib/api.ts surfaces it).
            return jsonify({"error": "Unauthorized"}), 401
        return f(*args, **kwargs)
    return decorated_function


# ─────────────────────────────────────────────────────────────
# ROUTES
# ─────────────────────────────────────────────────────────────

@app.route("/api/config")
@login_required
def get_config():
    """Expose feature flags to frontend JS."""
    return jsonify(FEATURE_FLAGS)


@app.route("/logout")
def logout():
    """Clear session and logout."""
    session.clear()
    accept = request.headers.get("Accept", "")
    if request.is_json or "application/json" in accept:
        return jsonify({"ok": True})
    # Direct browser hit (e.g. user typing /logout in URL bar): send them
    # to the Next.js login page, which the catch-all proxy now serves.
    return redirect("/login")


# ─────────────────────────────────────────────────────────────
# JSON AUTH ENDPOINTS (for the Next.js frontend)
# The HTML /login form is preserved above. /api/login is the JSON
# companion used by the Next client; both share the same Flask session.
# ─────────────────────────────────────────────────────────────

@app.route("/api/login", methods=["POST"])
def api_login():
    """JSON login endpoint. Body: {password}. Sets the Flask session cookie on success."""
    data = request.get_json(silent=True) or {}
    password = data.get("password") or ""
    expected_password = os.environ.get("APP_PASSWORD", "RRAM2026")
    if password != expected_password:
        return jsonify({"ok": False, "error": "Invalid access code"}), 401
    session["authenticated"] = True
    return jsonify({"ok": True})


@app.route("/api/me", methods=["GET"])
def api_me():
    """Return the auth state for the current session. Always 200 — the client uses
    `authenticated` to decide whether to redirect to /login."""
    return jsonify({"authenticated": bool(session.get("authenticated"))})


@app.route("/api/health", methods=["GET"])
def api_health():
    """Lightweight liveness probe. No auth required."""
    return jsonify({"ok": True, "service": "rr-powerbi-api"})


# ─────────────────────────────────────────────────────────────
# Frontend serving — Next.js static export
#
# Single-service deployment: `pnpm build` inside NewFrontEndToBePorted/
# produces ./out/, and Flask serves those files directly. No reverse
# proxy, no second Render service. /beta URLs are kept as 301 redirects
# to root so old bookmarks survive.
# ─────────────────────────────────────────────────────────────

NEXT_OUT_DIR = os.path.abspath(os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "NewFrontEndToBePorted",
    "out",
))

# Hand-authored standalone HTML pages served at their own routes (kept OUTSIDE
# the Next export's out/, which `pnpm build` wipes). e.g. /holidays.
STATIC_PAGES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static_pages")

# File extensions that are real assets — if missing, 404. (For non-asset
# paths we fall through to the SPA root index.html so client-side
# navigation works for any unknown route.)
_ASSET_EXTS = {
    ".css", ".js", ".mjs", ".map", ".json", ".txt", ".xml",
    ".ico", ".png", ".svg", ".jpg", ".jpeg", ".gif", ".webp",
    ".woff", ".woff2", ".ttf", ".otf", ".eot",
}


def _serve_static_export(subpath: str):
    """Resolve ``subpath`` against the Next.js export and stream the file back.

    Resolution order:
      1. Exact file (e.g. ``_next/static/foo.js``, ``icon.svg``)
      2. ``<path>.html``                 (e.g. /login → login.html)
      3. ``<path>/index.html``           (Next.js default route shape)
      4. Asset extension and not found → 404
      5. Anything else (i.e. an unknown HTML route) → root ``index.html``
         so the client-side router can take it from there.
    """
    if not os.path.isdir(NEXT_OUT_DIR):
        return jsonify({
            "error": "Frontend build not found. Run `pnpm build` inside "
                     "NewFrontEndToBePorted/ during the deploy build step.",
        }), 503

    safe = (subpath or "").lstrip("/")

    # 1. Exact file
    if safe:
        full = os.path.join(NEXT_OUT_DIR, safe)
        if os.path.isfile(full):
            return send_from_directory(NEXT_OUT_DIR, safe)

        # 2. <path>.html
        if not safe.endswith(".html"):
            html_path = safe + ".html"
            if os.path.isfile(os.path.join(NEXT_OUT_DIR, html_path)):
                return send_from_directory(NEXT_OUT_DIR, html_path)

        # 3. <path>/index.html
        idx_path = os.path.join(safe, "index.html")
        if os.path.isfile(os.path.join(NEXT_OUT_DIR, idx_path)):
            return send_from_directory(NEXT_OUT_DIR, idx_path)

        # 4. Asset extension but missing → real 404
        ext = os.path.splitext(safe)[1].lower()
        if ext in _ASSET_EXTS:
            return jsonify({"error": "Not Found", "path": safe}), 404

    # 5. SPA fallback (and the bare "/" route)
    return send_from_directory(NEXT_OUT_DIR, "index.html")


# /beta legacy redirects — preserve bookmarks from the previous frontend
# layout. 301 because the move is permanent.
@app.route("/beta", methods=["GET", "HEAD"])
@app.route("/beta/", methods=["GET", "HEAD"])
def beta_legacy_root():
    return redirect("/", code=301)


@app.route("/beta/<path:subpath>", methods=["GET", "HEAD"])
def beta_legacy_subpath(subpath):
    query = ("?" + request.query_string.decode("utf-8", "ignore")) if request.query_string else ""
    return redirect(f"/{subpath}{query}", code=301)


# Standalone hand-authored HTML page (public, no login). Literal route so it
# takes precedence over the catch-all SPA serve below.
@app.route("/holidays", methods=["GET", "HEAD"])
@app.route("/holidays/", methods=["GET", "HEAD"])
def holidays_page():
    return send_from_directory(STATIC_PAGES_DIR, "ME_A_Holidays_2026.html")


# ── Holidays editor API (powers the public /holidays page) ──────────────────
# Persistence: Cloudflare R2 (object holidays_2026.json) is the primary store;
# Postgres app_kv is a fallback if R2 isn't configured. Either survives restarts
# and is shared across visitors.
HOLIDAYS_R2_KEY = "app/holidays_2026.json"
HOLIDAYS_KV_KEY = "holidays_2026"
HOLIDAYS_EDIT_PASSWORD = os.environ.get("HOLIDAYS_EDIT_PASSWORD", "RRAM2026")
_HOLIDAYS_DEFAULT_PATH = os.path.join(STATIC_PAGES_DIR, "holidays_2026_default.json")


def _holidays_load_raw():
    """Load the saved holidays JSON string: R2 first, then Postgres KV."""
    try:
        raw = r2_get_text(HOLIDAYS_R2_KEY)
    except Exception:
        raw = None
    if raw:
        return raw
    try:
        return kv_get(HOLIDAYS_KV_KEY)
    except Exception:
        return None


def _holidays_save_raw(text):
    """Persist the holidays JSON string: try R2, then Postgres KV."""
    try:
        if r2_put_text(HOLIDAYS_R2_KEY, text):
            return True
    except Exception:
        pass
    try:
        return kv_set(HOLIDAYS_KV_KEY, text)
    except Exception:
        return False


def _holidays_default():
    try:
        with open(_HOLIDAYS_DEFAULT_PATH, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _sanitize_holidays(data):
    """Coerce posted data into {market: {holidays:[{date,name,note}], notes:[...]}}."""
    if not isinstance(data, dict):
        raise ValueError("data must be an object keyed by market")
    out = {}
    for market, rec in list(data.items())[:50]:
        market = str(market).strip()[:60]
        if not market or not isinstance(rec, dict):
            continue
        hols = []
        for h in (rec.get("holidays") or [])[:300]:
            if not isinstance(h, dict):
                continue
            date = str(h.get("date") or "").strip()[:10]
            name = str(h.get("name") or "").strip()[:140]
            note = h.get("note")
            note = (str(note).strip()[:300] or None) if note not in (None, "") else None
            if date and name:
                hols.append({"date": date, "name": name, "note": note})
        hols.sort(key=lambda x: x["date"])
        notes = [str(n).strip()[:300] for n in (rec.get("notes") or [])[:20] if str(n).strip()]
        weekend = rec.get("weekend") if rec.get("weekend") in ("fri-sat", "sat-sun") else "sat-sun"
        out[market] = {"weekend": weekend, "holidays": hols, "notes": notes}
    if not out:
        raise ValueError("no valid markets in data")
    return out


@app.route("/api/holidays", methods=["GET"])
def holidays_get():
    """Public — current holidays data (saved override if present, else bundled default)."""
    raw = _holidays_load_raw()
    payload, source = None, "default"
    if raw:
        try:
            payload = json.loads(raw)
            source = "db"
        except Exception:
            payload = None
    if payload is None:
        payload = _holidays_default()
    # Manual dump so market insertion order is preserved (jsonify sorts keys).
    return app.response_class(
        json.dumps({"data": payload, "source": source}, ensure_ascii=False),
        mimetype="application/json")


@app.route("/api/holidays/verify", methods=["POST"])
def holidays_verify():
    """Check the edit password (so the page can unlock edit mode)."""
    data = request.get_json(silent=True) or {}
    return jsonify({"ok": (data.get("password") or "") == HOLIDAYS_EDIT_PASSWORD})


@app.route("/api/holidays", methods=["POST"])
def holidays_save():
    """Persist edited holidays. Requires the edit password."""
    data = request.get_json(silent=True) or {}
    if (data.get("password") or "") != HOLIDAYS_EDIT_PASSWORD:
        return jsonify({"ok": False, "error": "Invalid edit password"}), 401
    try:
        clean = _sanitize_holidays(data.get("data"))
    except Exception as e:
        return jsonify({"ok": False, "error": f"Invalid data: {e}"}), 400
    if not _holidays_save_raw(json.dumps(clean, ensure_ascii=False)):
        return jsonify({"ok": False, "error": "Could not persist (storage unavailable)."}), 503
    return jsonify({"ok": True, "markets": len(clean)})


# Catch-all front-end serve. GET/HEAD only — POSTs are reserved for the
# explicit /api/* and /logout routes above. Flask routes by specificity,
# so this only fires for paths nothing else handled.
@app.route("/", defaults={"subpath": ""}, methods=["GET", "HEAD"])
@app.route("/<path:subpath>", methods=["GET", "HEAD"])
def serve_frontend(subpath):
    return _serve_static_export(subpath)


@app.route("/api/upload", methods=["POST"])
@login_required
def upload_files():
    """Accept one or more .xlsx files (Multipart or Base64 JSON), parse them, return JSON data."""
    files_to_process = []
    
    # Check for JSON Base64 upload (NetSkope Bypass)
    if request.is_json:
        data = request.get_json()
        if "files" in data:
            for f in data["files"]:
                fname = f.get("name")
                fdata = f.get("data") # data:application/vnd...;base64,....
                if fname and fdata:
                    try:
                        # Strip header if present
                        if "," in fdata:
                            _, b64data = fdata.split(",", 1)
                        else:
                            b64data = fdata
                        file_bytes = base64.b64decode(b64data)
                        files_to_process.append({"filename": fname, "bytes": file_bytes})
                    except Exception:
                        pass # Skip malformed

    # Fallback to standard Multipart upload
    if not files_to_process and "files" in request.files:
         for f in request.files.getlist("files"):
            if f.filename:
                files_to_process.append({"filename": f.filename, "bytes": f.read()})

    if not files_to_process:
        return jsonify({"error": "No files provided"}), 400

    sid = _get_session_id()
    results = {}
    errors = []

    for f in files_to_process:
        fname = f["filename"]
        file_bytes = f["bytes"]
        lower_fname = fname.lower()
        
        # Save to Database (New Feature)
        try:
            save_file_to_db(fname, file_bytes, sid)
        except Exception as e:
            print(f"Failed to save {fname} to DB: {e}")
        
        # ─── EXCEL (Universal Parser) ───
        if lower_fname.endswith((".xlsx", ".xls", ".xlsb", ".xlsm")):
            try:
                buf = io.BytesIO(file_bytes)
                parsed = parse_file(buf, filename=fname)
                # Sanitize for JSON (handle NaN, datetime, numpy, pandas types)
                parsed = _sanitize_for_json(parsed)
                results[fname] = parsed
                print(f"  Parsed {fname}: file_type={parsed.get('file_type', '??')}")

                # Store raw parsed data
                if sid not in _parsed_store:
                    _parsed_store[sid] = {}
                _parsed_store[sid][fname] = {
                    "type": "excel",
                    "file_type": parsed.get("file_type", "UNKNOWN"),
                    "parsed": parsed,
                    "file_bytes": file_bytes,
                }
            except Exception as e:
                import traceback
                traceback.print_exc()
                errors.append({"file": fname, "error": str(e)})

        # ─── PDF ───
        elif lower_fname.endswith(".pdf"):
            try:
                import pypdf
                buf = io.BytesIO(file_bytes)
                reader = pypdf.PdfReader(buf)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                
                if sid not in _parsed_store:
                    _parsed_store[sid] = {}
                _parsed_store[sid][fname] = {
                    "type": "pdf",
                    "text": text,
                    "file_bytes": file_bytes
                }
                results[fname] = {"text_preview": text[:200] + "..."}
            except Exception as e:
                errors.append({"file": fname, "error": f"PDF Error: {str(e)}"})

        # ─── WORD DOC ───
        elif lower_fname.endswith(".docx"):
            try:
                import docx
                buf = io.BytesIO(file_bytes)
                doc = docx.Document(buf)
                text = "\n".join([para.text for para in doc.paragraphs])
                
                if sid not in _parsed_store:
                    _parsed_store[sid] = {}
                _parsed_store[sid][fname] = {
                    "type": "docx",
                    "text": text,
                    "file_bytes": file_bytes
                }
                results[fname] = {"text_preview": text[:200] + "..."}
            except Exception as e:
                errors.append({"file": fname, "error": f"Docx Error: {str(e)}"})

        # ─── POWERPOINT ───
        elif lower_fname.endswith(".pptx"):
            try:
                from pptx import Presentation
                buf = io.BytesIO(file_bytes)
                prs = Presentation(buf)
                
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    slide_texts.append(para_text)
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                                if row_text.strip(" |"):
                                    slide_texts.append(row_text)
                    if slide_texts:
                        text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
                
                text = "\n\n".join(text_parts)
                
                if sid not in _parsed_store:
                    _parsed_store[sid] = {}
                _parsed_store[sid][fname] = {
                    "type": "pptx",
                    "text": text,
                    "file_bytes": file_bytes
                }
                results[fname] = {"text_preview": text[:200] + "..."}
            except Exception as e:
                errors.append({"file": fname, "error": f"PPTX Error: {str(e)}"})

        # ─── IMAGES ───
        elif lower_fname.endswith((".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif", ".gif")):
            try:
                # Convert to base64 for AI usage
                b64_img = base64.b64encode(file_bytes).decode('utf-8')
                
                if sid not in _parsed_store:
                    _parsed_store[sid] = {}
                _parsed_store[sid][fname] = {
                    "type": "image",
                    "base64": b64_img,
                    "mime": "image/png" if lower_fname.endswith(".png") else "image/jpeg",
                    "file_bytes": file_bytes
                }
                results[fname] = {"status": "Image stored for AI analysis"}
            except Exception as e:
                errors.append({"file": fname, "error": f"Image Error: {str(e)}"})

        else:
             errors.append({"file": fname, "error": "Unsupported file type."})
             continue

    if not results and errors:
        return jsonify({"error": "All files failed to process", "details": errors}), 400

    return jsonify({
        "files": results,
        "errors": errors if errors else None,
    })


@app.route("/api/export-pdf", methods=["POST"])
@login_required
def export_pdf():
    """Generate a PDF report from the uploaded data.

    Accepts two body shapes:

    1. New Next.js frontend (export-modal.tsx):
       {
         filename: "Foo.xlsx",
         file_type: "GLOBAL_HOPPER",
         format: "pdf",
         sections: { summary: true, charts: true, tables: true, insights: true }
       }

    2. Legacy clients:
       {
         selected_files: ["Foo.xlsx", ...],
         file_type: "...",
         sections_to_include: ["summary", "charts", ...],
         filters: {...},
         currency_symbol: "USD"
       }

    The handler normalises both into the legacy variables (selected_files,
    sections_to_include, filters, currency_symbol) before dispatching to
    the right `pdf_export.*` generator.
    """
    sid = _get_session_id()
    stored = _parsed_store.get(sid, {})

    if not stored:
        return jsonify({"error": "No data available. Please upload files first."}), 400

    data = request.get_json(silent=True) or {}
    currency_symbol = data.get("currency_symbol", "USD")
    file_type = data.get("file_type")
    filters = data.get("filters", {})
    # "detailed" -> long-form report with every analysis (charts + tables).
    # "ultra"    -> detailed report plus the ultra-detailed appendix (full
    #               register + VP/region breakdowns). Implies detailed.
    ultra = bool(data.get("ultra"))
    detailed = bool(data.get("detailed")) or ultra

    # -- selected_files (list) OR filename (single) ---------------------------
    selected_files = data.get("selected_files")
    if not selected_files:
        single = data.get("filename")
        if single:
            selected_files = [single]
        else:
            selected_files = list(stored.keys())

    # The per-section selection UI was removed, so the report is always the
    # full V8 report. We pass the full superset of section names both
    # generators understand (each ignores names it doesn't recognise) rather
    # than None — the V8 Opportunity-Tracker generator does ``"kpis" in
    # sections_to_include`` and would raise on None.
    sections_to_include = [
        # Hopper section names
        "summary", "charts", "customer_analysis", "engine_analysis",
        "pipeline", "restructure", "top_opportunities", "customer_breakdown",
        # Opportunity Tracker section names
        "kpis", "top_opps", "estimation_level", "opps_threats",
        "project_summary", "timeline",
    ]

    # Resolve the active file (tolerate filename not matching any stored key).
    first_key = selected_files[0] if selected_files else list(stored.keys())[0]
    if first_key not in stored:
        first_key = list(stored.keys())[0]
    first_parsed = stored[first_key]["parsed"]

    if file_type == 'GLOBAL_HOPPER' or first_parsed.get("file_type") == 'GLOBAL_HOPPER':
        try:
            if detailed:
                from pdf_export import generate_hopper_detailed_pdf_report
                pdf_bytes = generate_hopper_detailed_pdf_report(
                    parsed_data=first_parsed,
                    sections_to_include=sections_to_include,
                    filters=filters,
                    ultra=ultra,
                )
            else:
                from pdf_export import generate_hopper_pdf_report
                pdf_bytes = generate_hopper_pdf_report(
                    parsed_data=first_parsed,
                    sections_to_include=sections_to_include,
                    filters=filters
                )
            filename = ("Global_Hopper_Ultra_Detailed_Report.pdf" if ultra
                        else "Global_Hopper_Detailed_Report.pdf" if detailed
                        else "Global_Hopper_Report.pdf")
            return send_file(
                io.BytesIO(pdf_bytes),
                mimetype="application/pdf",
                as_attachment=True,
                download_name=filename,
            )
        except Exception as e:
            import traceback
            traceback.print_exc()
            return jsonify({"error": f"Global Hopper PDF generation failed: {str(e)}"}), 500

    if file_type == 'OPPORTUNITY_TRACKER' or first_parsed.get("file_type") == 'OPPORTUNITY_TRACKER':
        try:
            from pdf_export import generate_opp_pdf_report
            pdf_bytes = generate_opp_pdf_report(
                parsed_data=first_parsed,
                sections_to_include=sections_to_include,
                filters=filters
            )
            cust_name = first_parsed.get("metadata", {}).get("customer", "Report").replace(" ", "_")
            filename = f"Opportunity_Tracker_{cust_name}.pdf"
            return send_file(
                io.BytesIO(pdf_bytes),
                mimetype="application/pdf",
                as_attachment=True,
                download_name=filename,
            )
        except Exception as e:
            import traceback
            traceback.print_exc()
            return jsonify({"error": f"Opp Tracker PDF generation failed: {str(e)}"}), 500

    # Fallback to SOA export. The SOA parser returns ``sections`` as a list of
    # {name, section_type, items[], total, overdue} dicts (each item carries
    # lowercase keys: reference/doc_date/due_date/amount/days_late/...), so we
    # normalise here into the (metadata, grand_totals, filtered_df,
    # sections_summary) shape the generator expects.
    metadata = first_parsed.get("metadata", {})
    grand_totals = dict(first_parsed.get("grand_totals", {}) or {})

    raw_sections = first_parsed.get("sections", [])
    if isinstance(raw_sections, dict):
        raw_sections = [{"name": k, **(v or {})} for k, v in raw_sections.items()]

    sections_summary = {}
    normalised_items = []
    total_charges = 0.0
    item_count = 0
    for sec in (raw_sections or []):
        name = sec.get("name", "Section")
        sec_items = sec.get("items") or sec.get("rows") or []
        charges = sum(_v for it in sec_items if (_v := (it.get("amount") or 0)) > 0)
        credits = sum(_v for it in sec_items if (_v := (it.get("amount") or 0)) < 0)
        total_charges += charges
        item_count += len(sec_items)
        sections_summary[name] = {
            "total": sec.get("total") if sec.get("total") is not None else (charges + credits),
            "charges": charges,
            "credits": credits,
            "overdue": sec.get("overdue") or 0,
            "items": len(sec_items),
        }
        for it in sec_items:
            normalised_items.append({
                "Section": name,
                "Reference": it.get("reference"),
                "Document Date": it.get("doc_date"),
                "Due Date": it.get("due_date"),
                "Amount": it.get("amount"),
                "Status": it.get("customer_comments") or it.get("rr_comments"),
                "Entry Type": sec.get("section_type") or it.get("text"),
                "Days Late": it.get("days_late"),
            })

    grand_totals.setdefault("total_charges", total_charges)
    grand_totals.setdefault("item_count", item_count)
    filtered_df = pd.DataFrame(normalised_items)
    source_files = selected_files if len(selected_files) > 1 else None

    try:
        from pdf_export import generate_pdf_report
        pdf_bytes = generate_pdf_report(
            metadata=metadata,
            grand_totals=grand_totals,
            filtered_df=filtered_df,
            sections_summary=sections_summary,
            source_files=source_files,
            currency_symbol=currency_symbol,
        )

        cust_name = metadata.get("customer_name", "Report").replace(" ", "_")
        filename = f"SOA_Report_{cust_name}.pdf"

        return send_file(
            io.BytesIO(pdf_bytes),
            mimetype="application/pdf",
            as_attachment=True,
            download_name=filename,
        )
    except Exception as e:
        return jsonify({"error": f"PDF generation failed: {str(e)}"}), 500


# ─────────────────────────────────────────────────────────────
# AI REPORT (Kimi K2.6) — background job + polling
# ─────────────────────────────────────────────────────────────

@app.route("/api/ai-report", methods=["POST"])
@login_required
def ai_report_start():
    """Kick off an AI-designed PDF report. Returns a job_id immediately; the
    work runs in a background thread (poll /api/ai-report/<id>)."""
    _prune_ai_jobs()
    sid = _get_session_id()
    stored = _parsed_store.get(sid, {})
    if not stored:
        return jsonify({"error": "No data available. Please upload files first."}), 400

    data = request.get_json(silent=True) or {}
    mode = data.get("mode", "catalog")
    if mode not in ("catalog", "charts", "html"):
        mode = "catalog"
    provider = data.get("provider", "nvidia")
    if provider not in ("nvidia", "aistudio"):
        provider = "nvidia"
    filters = data.get("filters", {}) or {}

    selected = data.get("filename")
    keys = list(stored.keys())
    first_key = selected if selected in stored else (keys[0] if keys else None)
    if not first_key:
        return jsonify({"error": "No data available."}), 400
    first_parsed = stored[first_key]["parsed"]
    ftype = data.get("file_type") or first_parsed.get("file_type")
    if ftype != "GLOBAL_HOPPER":
        return jsonify({"error": "AI reports currently support Global Hopper files only."}), 400

    job_id = uuid.uuid4().hex
    with _ai_report_lock:
        _ai_report_jobs[job_id] = {
            "status": "queued", "mode": mode, "provider": provider, "progress": "Queued…",
            "pdf": None, "filename": None, "error": None, "note": None,
            "created": time.time(),
        }

    def _run(parsed_snapshot, flt, m, prov):
        with _ai_report_lock:
            if job_id in _ai_report_jobs:
                _ai_report_jobs[job_id]["status"] = "running"

        def progress(msg):
            with _ai_report_lock:
                if job_id in _ai_report_jobs:
                    _ai_report_jobs[job_id]["progress"] = msg

        try:
            from ai_report import generate_ai_report
            pdf_bytes, filename, note = generate_ai_report(parsed_snapshot, flt, m, progress, prov)
            with _ai_report_lock:
                if job_id in _ai_report_jobs:
                    _ai_report_jobs[job_id].update(status="done", pdf=pdf_bytes,
                                                   filename=filename, note=note, progress="Done")
        except Exception as e:
            import traceback
            traceback.print_exc()
            with _ai_report_lock:
                if job_id in _ai_report_jobs:
                    _ai_report_jobs[job_id].update(status="failed", error=str(e), progress="Failed")

    threading.Thread(target=_run, args=(first_parsed, filters, mode, provider), daemon=True).start()
    return jsonify({"job_id": job_id, "mode": mode, "provider": provider})


@app.route("/api/ai-report/<job_id>", methods=["GET"])
@login_required
def ai_report_status(job_id):
    with _ai_report_lock:
        job = _ai_report_jobs.get(job_id)
        if not job:
            return jsonify({"error": "Unknown job"}), 404
        return jsonify({"status": job["status"], "progress": job["progress"],
                        "error": job["error"], "note": job["note"],
                        "filename": job["filename"], "mode": job["mode"],
                        "provider": job.get("provider", "nvidia")})


@app.route("/api/ai-report/<job_id>/download", methods=["GET"])
@login_required
def ai_report_download(job_id):
    with _ai_report_lock:
        job = _ai_report_jobs.get(job_id)
        if not job:
            return jsonify({"error": "Unknown job"}), 404
        if job["status"] != "done" or not job["pdf"]:
            return jsonify({"error": "Report not ready"}), 409
        pdf_bytes, filename = job["pdf"], (job["filename"] or "AI_Report.pdf")
        _ai_report_jobs.pop(job_id, None)   # evict after download
    return send_file(io.BytesIO(pdf_bytes), mimetype="application/pdf",
                     as_attachment=True, download_name=filename)


# ─────────────────────────────────────────────────────────────
# FILE MANAGEMENT ROUTES
# ─────────────────────────────────────────────────────────────

@app.route("/api/files", methods=["GET"])
@login_required
def list_files():
    """List all uploaded files from the database."""
    from storage import get_all_files
    files = get_all_files()
    return jsonify(files)


@app.route("/api/files/<int:file_id>", methods=["GET"])
@login_required
def download_file(file_id):
    """Download a specific file from the database."""
    from storage import get_file_by_id
    file_record = get_file_by_id(file_id)
    
    if not file_record:
        return jsonify({"error": "File not found"}), 404
        
    filename = file_record['filename']
    file_data = file_record['file_data'] # This is bytes (BYTEA)
    
    # BytesIO for sending file
    return send_file(
        io.BytesIO(file_data),
        as_attachment=True,
        download_name=filename
    )


@app.route("/api/files/upload", methods=["POST"])
@login_required
def upload_files_to_db():
    """Upload files directly to the database without parsing."""
    try:
        data = request.json
        if not data or 'files' not in data:
            return jsonify({"error": "No files provided"}), 400

        files = data['files']
        saved_count = 0
        
        for f in files:
            fname = f.get("name")
            file_data_b64 = f.get("data") # "data:application/pdf;base64,....."
            
            if not fname or not file_data_b64:
                continue
                
            # Decode base64
            if "," in file_data_b64:
                _, b64_str = file_data_b64.split(",", 1)
            else:
                b64_str = file_data_b64
                
            file_bytes = base64.b64decode(b64_str)
            
            # Save to DB
            save_file_to_db(fname, file_bytes)
            saved_count += 1
            
        return jsonify({"message": f"Successfully uploaded {saved_count} files", "count": saved_count})

    except Exception as e:
        print(f"Error uploading files: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/files/<int:file_id>", methods=["DELETE"])
@login_required
def delete_file(file_id):
    """Delete a specific file from the database."""
    from storage import delete_file_by_id
    success = delete_file_by_id(file_id)

    if success:
        return jsonify({"message": "File deleted successfully"})
    else:
        return jsonify({"error": "Failed to delete file"}), 500


@app.route("/api/parsed/<path:fname>", methods=["DELETE"])
@login_required
def delete_parsed_file(fname):
    """Remove a parsed file from the in-memory session store (dashboard-side removal)."""
    sid = _get_session_id()
    if sid in _parsed_store and fname in _parsed_store[sid]:
        del _parsed_store[sid][fname]
        return jsonify({"ok": True, "filename": fname})
    return jsonify({"ok": False, "error": "not found"}), 404


# ─────────────────────────────────────────────────────────────
# R2 CLOUD STORAGE ROUTES (V2 Upload)
# ─────────────────────────────────────────────────────────────

@app.route("/api/r2/chunk-init", methods=["POST"])
@login_required
def r2_chunk_init():
    """Start an S3 multipart upload on R2. Returns session_id for subsequent chunk uploads."""
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")
    total_chunks = data.get("total_chunks")

    if not filename or not total_chunks:
        return jsonify({"error": "filename and total_chunks required"}), 400

    # Generate R2 key and start multipart upload
    r2_key = generate_r2_key(filename)
    r2_upload_id = create_multipart_upload(r2_key)
    if not r2_upload_id:
        return jsonify({"error": "Failed to start multipart upload on R2"}), 500

    session_id = uuid.uuid4().hex
    _multipart_sessions[session_id] = {
        "filename": filename,
        "r2_key": r2_key,
        "r2_upload_id": r2_upload_id,
        "total": int(total_chunks),
        "parts": [],       # [{"PartNumber": int, "ETag": str}, ...]
        "file_size": 0,
        "created": time.time(),
    }
    print(f"R2 multipart upload initialized: {session_id} for {filename} ({total_chunks} chunks)")
    return jsonify({"upload_id": session_id})


@app.route("/api/r2/chunk-upload", methods=["POST"])
@login_required
def r2_chunk_upload():
    """Receive a single chunk, decode Base64, stream directly to R2 as a part. No in-memory accumulation."""
    data = request.get_json(silent=True) or {}
    session_id = data.get("upload_id")
    chunk_index = data.get("chunk_index")
    chunk_data = data.get("data")  # raw Base64 string

    if not session_id or chunk_index is None or not chunk_data:
        return jsonify({"error": "upload_id, chunk_index, and data required"}), 400

    sess = _multipart_sessions.get(session_id)
    if not sess:
        return jsonify({"error": "Invalid or expired upload_id"}), 404

    try:
        decoded = base64.b64decode(chunk_data)
    except Exception:
        return jsonify({"error": "Invalid base64 data"}), 400

    # S3 part numbers are 1-based
    part_number = int(chunk_index) + 1

    # Upload part directly to R2 (only ~8MB in memory at a time)
    etag = upload_part(sess["r2_key"], sess["r2_upload_id"], part_number, decoded)
    if not etag:
        return jsonify({"error": f"Failed to upload part {part_number} to R2"}), 500

    sess["parts"].append({"PartNumber": part_number, "ETag": etag})
    sess["file_size"] += len(decoded)

    # decoded is now garbage-collectable
    del decoded

    received = len(sess["parts"])
    total = sess["total"]
    print(f"  Part {part_number}/{total} streamed to R2 for {sess['filename']} ({len(chunk_data)} b64 chars)")

    return jsonify({"received": received, "total": total})


@app.route("/api/r2/chunk-finalize", methods=["POST"])
@login_required
def r2_chunk_finalize():
    """Complete the multipart upload on R2 (R2 assembles the parts server-side). Save metadata to DB."""
    data = request.get_json(silent=True) or {}
    session_id = data.get("upload_id")

    if not session_id:
        return jsonify({"error": "upload_id required"}), 400

    sess = _multipart_sessions.get(session_id)
    if not sess:
        return jsonify({"error": "Invalid or expired upload_id"}), 404

    total = sess["total"]
    if len(sess["parts"]) != total:
        # Abort and clean up
        abort_multipart_upload(sess["r2_key"], sess["r2_upload_id"])
        del _multipart_sessions[session_id]
        return jsonify({
            "error": f"Missing parts: received {len(sess['parts'])}/{total}"
        }), 400

    # Sort parts by PartNumber (required by S3)
    parts_sorted = sorted(sess["parts"], key=lambda p: p["PartNumber"])

    # Complete — R2 assembles the file server-side, zero memory usage here
    success = complete_multipart_upload(sess["r2_key"], sess["r2_upload_id"], parts_sorted)
    if not success:
        abort_multipart_upload(sess["r2_key"], sess["r2_upload_id"])
        del _multipart_sessions[session_id]
        return jsonify({"error": "Failed to complete multipart upload on R2"}), 500

    filename = sess["filename"]
    r2_key = sess["r2_key"]
    file_size = sess["file_size"]
    public_url = f"{R2_PUBLIC_URL}/{r2_key}" if R2_PUBLIC_URL else None

    # Clean up session
    del _multipart_sessions[session_id]

    # Save metadata to PostgreSQL (no file bytes — just a pointer to R2)
    from storage import save_r2_file_metadata
    sid = _get_session_id()
    row_id = save_r2_file_metadata(
        filename=filename,
        r2_key=r2_key,
        public_url=public_url,
        file_size=file_size,
        session_id=sid,
    )

    return jsonify({
        "id": row_id,
        "filename": filename,
        "r2_key": r2_key,
        "public_url": public_url,
        "file_size": file_size,
    })


@app.route("/api/r2/files", methods=["GET"])
@login_required
def list_r2_files():
    """List all R2-stored files (metadata from DB)."""
    from storage import get_all_r2_files
    files = get_all_r2_files()
    return jsonify(files)


@app.route("/api/r2/files/<int:file_id>", methods=["GET"])
@login_required
def download_r2_file(file_id):
    """Download a file from R2 by DB id."""
    from storage import get_r2_file_by_id
    meta = get_r2_file_by_id(file_id)
    if not meta:
        return jsonify({"error": "File not found"}), 404

    file_bytes = download_from_r2(meta["r2_key"])
    if not file_bytes:
        return jsonify({"error": "Failed to download from R2"}), 500

    return send_file(
        io.BytesIO(file_bytes),
        as_attachment=True,
        download_name=meta["filename"],
    )


@app.route("/api/r2/files/<int:file_id>", methods=["DELETE"])
@login_required
def delete_r2_file(file_id):
    """Delete a file from R2 and its metadata from DB."""
    from storage import delete_r2_file_by_id
    r2_key = delete_r2_file_by_id(file_id)
    if not r2_key:
        return jsonify({"error": "File not found"}), 404

    # Delete from R2
    delete_from_r2(r2_key)
    return jsonify({"message": "File deleted from R2 and database"})


@app.route("/api/r2/files/<int:file_id>/parse", methods=["POST"])
@login_required
def parse_r2_file(file_id):
    """Download file from R2, parse it, return JSON + store in memory for dashboard."""
    from storage import get_r2_file_by_id
    meta = get_r2_file_by_id(file_id)
    if not meta:
        return jsonify({"error": "File not found"}), 404

    file_bytes = download_from_r2(meta["r2_key"])
    if not file_bytes:
        return jsonify({"error": "Failed to download from R2"}), 500

    fname = meta["filename"]
    sid = _get_session_id()

    try:
        buf = io.BytesIO(file_bytes)
        parsed = parse_file(buf, filename=fname)
        parsed = _sanitize_for_json(parsed)

        # Store in memory for dashboard use
        if sid not in _parsed_store:
            _parsed_store[sid] = {}
        _parsed_store[sid][fname] = {
            "type": "excel",
            "file_type": parsed.get("file_type", "UNKNOWN"),
            "parsed": parsed,
            "file_bytes": file_bytes,
        }

        return jsonify({"files": {fname: parsed}})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"Parse failed: {str(e)}"}), 500


# ─────────────────────────────────────────────────────────────
# AI CHAT
# ─────────────────────────────────────────────────────────────

@app.route("/api/chat", methods=["POST"])
@login_required
def chat():
    """Handle AI chat messages."""
    sid = _get_session_id()
    stored = _parsed_store.get(sid, {})

    if not stored:
        return jsonify({"error": "No data available. Please upload files first."}), 400

    data = request.get_json(silent=True) or {}
    user_message = data.get("message", "").strip()
    model_choice = data.get("model", None)

    if not user_message:
        return jsonify({"error": "No message provided"}), 400

    # Build serialized data dict for system prompt
    serialized_data = {}
    images_to_attach = []           # For OpenAI-format models (Qwen etc.)
    gemini_file_attachments = []    # For Gemini 3 Pro native multimodal

    for fname, fstore in stored.items():
        ftype = fstore.get("type", "excel")
        
        if ftype in ["pdf", "docx", "pptx"]:
            serialized_data[fname] = {"type": ftype, "text": fstore.get("text")}
            
            # For Gemini: attach the raw PDF natively (Gemini can read PDFs directly)
            if ftype == "pdf" and fstore.get("file_bytes"):
                gemini_file_attachments.append({
                    "mime_type": "application/pdf",
                    "base64": base64.b64encode(fstore["file_bytes"]).decode("utf-8"),
                    "filename": fname
                })
            
        elif ftype == "image":
            serialized_data[fname] = {"type": "image"}
            # OpenAI-format image for non-Gemini models
            images_to_attach.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:{fstore['mime']};base64,{fstore['base64']}"
                }
            })
            # Gemini-native image attachment
            gemini_file_attachments.append({
                "mime_type": fstore["mime"],
                "base64": fstore["base64"],
                "filename": fname
            })
        else:
            # Excel — use the parsed data in system prompt
            # Universal parser results are already JSON-serializable
            parsed = fstore.get("parsed", {})
            if serialize_parsed_data and isinstance(parsed, dict) and "sections" in parsed and hasattr(parsed.get("sections", {}), 'items'):
                try:
                    serialized_data[fname] = serialize_parsed_data(parsed)
                except Exception:
                    serialized_data[fname] = parsed
            else:
                serialized_data[fname] = parsed

    system_prompt = build_system_prompt(serialized_data)

    # Get/create chat history for this session
    if sid not in _chat_history:
        _chat_history[sid] = []

    # Construct User Message (Text + Images for OpenAI-format models)
    user_content = []
    if user_message:
        user_content.append({"type": "text", "text": user_message})
    user_content.extend(images_to_attach)

    # Add user message to history
    _chat_history[sid].append({"role": "user", "content": user_content})

    # Call AI — pass file_attachments for Gemini models
    result = call_openrouter(
        _chat_history[sid],
        system_prompt,
        model=model_choice,
        file_attachments=gemini_file_attachments if model_choice and model_choice.startswith("gemini/") else None
    )

    if result.get("error"):
        _chat_history[sid].pop()
        return jsonify({"error": result["error"]}), 502

    # Add assistant response to history
    _chat_history[sid].append({"role": "assistant", "content": result.get("content", "")})

    return jsonify({
        "content": result.get("content", ""),
        "charts": result.get("charts", []),
        "emails": result.get("emails", []),
    })


@app.route("/api/chat/clear", methods=["POST"])
@login_required
def clear_chat():
    """Clear chat history for current session."""
    sid = _get_session_id()
    _chat_history[sid] = []
    return jsonify({"status": "ok"})


@app.route("/api/compare", methods=["POST"])
@login_required
def compare_models():
    """Run the same prompt against all 3 models in parallel."""
    sid = _get_session_id()
    stored = _parsed_store.get(sid, {})

    if not stored:
        return jsonify({"error": "No data available. Please upload files first."}), 400

    data = request.get_json(silent=True) or {}
    user_message = data.get("message", "").strip()

    if not user_message:
        return jsonify({"error": "No message provided"}), 400

    # Build system prompt
    serialized_data = {}
    for fname, fstore in stored.items():
        serialized_data[fname] = serialize_parsed_data(fstore["parsed"])
    
    system_prompt = build_system_prompt(serialized_data)
    
    # Models to compare
    models = [
        {"id": "qwen/qwen3-vl-235b-a22b-thinking", "name": "Qwen 3 VL (OpenRouter)"},
        {"id": "digitalocean/openai-gpt-oss-120b", "name": "GPT 120b (DigitalOcean)"},
        {"id": "nvidia/moonshotai/kimi-k2.5", "name": "Kimi K2.5 (NVIDIA)"},
        {"id": "gemini/gemini-3-pro-preview", "name": "Gemini 3 Pro (Google)"},
    ]

    results = []

    def fetch_model_response(model_info):
        start_t = time.time()
        m_id = model_info["id"]
        m_name = model_info["name"]
        
        # Use a temporary history with just this message for the comparison
        temp_history = [{"role": "user", "content": user_message}]
        
        try:
            # Pass file attachments for Gemini models
            # We need to reconstruct gemini_file_attachments here similar to /api/chat
            # Use a helper or just inline simple reconstruction
            gemini_atts = []
            for fname, fstore in stored.items():
                if m_id.startswith("gemini/") and fstore.get("type") == "pdf" and fstore.get("file_bytes"):
                     gemini_atts.append({
                        "mime_type": "application/pdf",
                        "base64": base64.b64encode(fstore["file_bytes"]).decode("utf-8"),
                        "filename": fname
                    })
                elif m_id.startswith("gemini/") and fstore.get("type") == "image":
                     gemini_atts.append({
                        "mime_type": fstore["mime"],
                        "base64": fstore["base64"],
                        "filename": fname
                    })

            resp = call_openrouter(temp_history, system_prompt, model=m_id, file_attachments=gemini_atts)
            duration = time.time() - start_t
            
            return {
                "model_id": m_id,
                "model_name": m_name,
                "content": resp.get("content", ""),
                "error": resp.get("error"),
                "time": f"{duration:.2f}s"
            }
        except Exception as e:
            return {
                "model_id": m_id,
                "model_name": m_name,
                "content": None,
                "error": str(e),
                "time": f"{time.time() - start_t:.2f}s"
            }

    # Execute in parallel
    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as executor:
        future_to_model = {executor.submit(fetch_model_response, m): m for m in models}
        for future in concurrent.futures.as_completed(future_to_model):
            results.append(future.result())

    # Sort results to match input order for consistency in UI
    results.sort(key=lambda x: [m["id"] for m in models].index(x["model_id"]))

    return jsonify({"results": results})


# ─────────────────────────────────────────────────────────────
# RUN
# ─────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import webbrowser
    import threading

    port = int(os.environ.get("PORT", 5000))
    print(f"\n  Rolls-Royce Data Visualizer")
    print(f"  http://localhost:{port}\n")

    # Auto-open browser after short delay
    def open_browser():
        import time
        time.sleep(1.5)
        webbrowser.open(f"http://localhost:{port}")

    threading.Thread(target=open_browser, daemon=True).start()

    app.run(host="0.0.0.0", port=port, debug=True, use_reloader=False)
